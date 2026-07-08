"""Build Session 1 deck for Intro to AI (v3 — dark mode, video-dominant layout).

Run:  /tmp/pptx-venv/bin/python session-1/build_deck.py
Out:  session-1/deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---- Palette (dark mode) ----
BG = RGBColor(0x1A, 0x19, 0x17)
TEXT = RGBColor(0xF8, 0xF6, 0xF2)
MUTED = RGBColor(0xA8, 0xA3, 0x9E)
CARD_DARK = RGBColor(0x2A, 0x28, 0x25)
ACCENT = RGBColor(0x7B, 0xA8, 0xE8)     # bright blue
WHY = RGBColor(0xD4, 0x96, 0x58)        # amber  (Beat 2 / 7 ceremonial)
P1 = RGBColor(0x9B, 0x8F, 0xE8)         # purple (intelligence / Beat 1 / Beat 5)
P2 = RGBColor(0x5F, 0xC8, 0xA0)         # green  (problems AI solves / journal-automate)
P3 = RGBColor(0xE8, 0xA0, 0x62)         # ochre  (backwards walk)
HERO_FG = RGBColor(0xF8, 0xF6, 0xF2)
HERO_DIM = RGBColor(0xCC, 0xC8, 0xC2)

HEAD_FONT = "Georgia"
BODY_FONT = "Helvetica Neue"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---- Grid (consistent across every content slide) ----
MARGIN = Inches(0.6)
TITLE_Y = Inches(0.55)
TITLE_H = Inches(1.25)
CONTENT_Y = Inches(2.0)
CONTENT_H = Inches(4.6)
FOOTER_Y = Inches(7.0)

VIDEO_W = Inches(8.4)            # ~63% of slide width
VIDEO_X = MARGIN                 # 0.6"
SIDEBAR_X = Inches(9.2)
SIDEBAR_W = Inches(3.55)


# ---------- Core helpers ----------

def new_pres():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, x, y, w, h, *, font=BODY_FONT, size=18, bold=False,
             italic=False, color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_title(slide, text, size=44):
    add_text(slide, text, MARGIN, TITLE_Y, Inches(12.1), TITLE_H,
             font=HEAD_FONT, size=size, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)


def add_footer(slide):
    add_text(slide, "INTRO TO AI  ·  SESSION 1",
             MARGIN, FOOTER_Y, Inches(12.1), Inches(0.35),
             font=BODY_FONT, size=10, color=MUTED, align=PP_ALIGN.LEFT)


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


# ---------- Video placeholder (compact, video-dominant) ----------

def add_video_placeholder(slide, label, search_hint, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 VIDEO_X, CONTENT_Y, VIDEO_W, CONTENT_H)
    box.adjustments[0] = 0.025
    box.fill.solid(); box.fill.fore_color.rgb = CARD_DARK
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(30); tf.margin_right = Pt(30)
    tf.margin_top = Pt(20); tf.margin_bottom = Pt(20)

    # ▶ glyph, large
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "▶"
    r.font.name = BODY_FONT; r.font.size = Pt(64); r.font.color.rgb = color

    # Label
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)
    r2 = p2.add_run(); r2.text = label
    r2.font.name = HEAD_FONT; r2.font.size = Pt(22); r2.font.color.rgb = TEXT

    # Search hint
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(12)
    r3 = p3.add_run(); r3.text = f"YouTube search: “{search_hint}”"
    r3.font.name = BODY_FONT; r3.font.size = Pt(12); r3.font.italic = True
    r3.font.color.rgb = MUTED


# ---------- Sidebar ----------

def render_sidebar(slide, eyebrow, bullets, color, footnote=None):
    """Right-column context: small-caps eyebrow + accent rule + bullets."""
    y = CONTENT_Y

    # Eyebrow
    add_text(slide, eyebrow.upper(), SIDEBAR_X, y, SIDEBAR_W, Inches(0.35),
             font=BODY_FONT, size=11, bold=True, color=color, align=PP_ALIGN.LEFT)

    # Thin accent rule (a 0.06" tall colored rectangle, no border)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  SIDEBAR_X, y + Inches(0.45),
                                  Inches(0.6), Inches(0.025))
    rule.line.fill.background()
    rule.fill.solid(); rule.fill.fore_color.rgb = color
    rule.shadow.inherit = False

    # Bullets
    bullet_box = slide.shapes.add_textbox(
        SIDEBAR_X, y + Inches(0.65), SIDEBAR_W, CONTENT_H - Inches(0.65))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "•  " + b
        r.font.name = BODY_FONT; r.font.size = Pt(13); r.font.color.rgb = TEXT

    # Optional footnote in muted italic at bottom of sidebar
    if footnote:
        add_text(slide, footnote,
                 SIDEBAR_X, CONTENT_Y + CONTENT_H - Inches(0.55),
                 SIDEBAR_W, Inches(0.5),
                 font=BODY_FONT, size=10, italic=True, color=MUTED)


# ---------- Slide builders ----------

def slide_title(prs):
    s = add_blank(prs)
    # Top tag
    add_text(s, "INTRO TO AI  ·  SESSION 1",
             MARGIN, Inches(0.6), Inches(12.1), Inches(0.4),
             font=BODY_FONT, size=13, color=HERO_DIM, align=PP_ALIGN.CENTER)
    # Hero title
    add_text(s, "What even is AI?",
             MARGIN, Inches(2.4), Inches(12.1), Inches(2.4),
             font=HEAD_FONT, size=120, color=HERO_FG,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Just the date — no tagline
    add_text(s, "Tue, June 23",
             MARGIN, Inches(5.6), Inches(12.1), Inches(0.5),
             font=BODY_FONT, size=18, color=HERO_DIM, align=PP_ALIGN.CENTER)
    add_notes(s, (
        "~70–75 min total. Spine: humans have always been problem-solvers; AI is the newest "
        "chapter; the problems still unsolved are Shreya's to solve.\n\n"
        "Arc: wow → her wishes → backwards walk → break → what AI solves now → why copy "
        "humans → what AI still can't do → why bother learning → journal.\n\n"
        "Close on the unanswered question — DO NOT answer it."
    ))


def video_slide(prs, title, video_label, search_hint, sidebar_eyebrow, sidebar_bullets,
                color, notes, sidebar_footnote=None):
    s = add_blank(prs)
    add_title(s, title)
    add_video_placeholder(s, video_label, search_hint, color)
    render_sidebar(s, sidebar_eyebrow, sidebar_bullets, color, footnote=sidebar_footnote)
    add_footer(s)
    add_notes(s, notes)
    return s


def slide_beat2_wishes(prs):
    """Interactive prompt card — no video."""
    s = add_blank(prs)
    add_title(s, "What do you wish you didn't have to do?")

    def prompt_card(x, w, eyebrow, text, color):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, CONTENT_Y, w, CONTENT_H)
        c.adjustments[0] = 0.03
        c.fill.solid(); c.fill.fore_color.rgb = CARD_DARK
        c.line.color.rgb = color; c.line.width = Pt(1.5)
        c.shadow.inherit = False
        tf = c.text_frame
        tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(32); tf.margin_right = Pt(32)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = eyebrow.upper()
        r.font.name = BODY_FONT; r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = color
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(20)
        r2 = p2.add_run(); r2.text = text
        r2.font.name = HEAD_FONT; r2.font.size = Pt(26); r2.font.color.rgb = TEXT

    half = (SLIDE_W - MARGIN * 3) / 2
    prompt_card(MARGIN, half,
                "Ask — write each answer on a card",
                "“What's something\nyou HAVE to do\nbut wish you didn't?”",
                WHY)
    prompt_card(MARGIN * 2 + half, half,
                "Optional follow-up",
                "“What's something\nyou WISH you could do\nbut can't?”",
                P1)

    add_text(s,
             "Keep the cards visible — they come back at Beat 6 (AI still can't do it) and the Journal wrap.",
             MARGIN, Inches(6.6), Inches(12.1), Inches(0.4),
             font=BODY_FONT, size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s)
    add_notes(s, (
        "Interactive — NO video. Make it personal and instant.\n\n"
        "Ask: 'What's something you HAVE to do but wish you didn't?' She might say cleaning "
        "her room, folding laundry, practicing scales, homework. Write each answer on a "
        "physical card. Keep them visible all session.\n\n"
        "Optional follow-up (quietly seeds the whales in Beat 4): 'What's something you "
        "WISH you could do but can't?' — talking to animals, flying, being in two places "
        "at once.\n\n"
        "BRIDGE: every tool humans ever built started exactly here — someone wishing a "
        "hard thing were easier. That's the doorway into the backwards walk."
    ))


def slide_intelligence_open(prs):
    """5C — open question, no video."""
    s = add_blank(prs)
    add_title(s, "So… what IS intelligence?")
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              MARGIN, CONTENT_Y, Inches(12.1), CONTENT_H)
    card.adjustments[0] = 0.03
    card.fill.solid(); card.fill.fore_color.rgb = CARD_DARK
    card.line.color.rgb = P1; card.line.width = Pt(1.5)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(48); tf.margin_right = Pt(48)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Collect answers from Shreya AND the parents.\nWrite them on cards. Leave it open."
    r.font.name = HEAD_FONT; r.font.size = Pt(30); r.font.color.rgb = TEXT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(22)
    r2 = p2.add_run(); r2.text = "Even the experts argue about this."
    r2.font.name = BODY_FONT; r2.font.size = Pt(17); r2.font.italic = True
    r2.font.color.rgb = P1

    add_text(s, "Leaving it unresolved is the design, not a gap.",
             MARGIN, Inches(6.6), Inches(12.1), Inches(0.4),
             font=BODY_FONT, size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s)
    add_notes(s, (
        "Close Beat 5 with the question Shreya gets to keep. Collect answers from her AND "
        "the parents (especially fun with two PhD parents — they'll disagree, and that's "
        "the point). Write them on cards. Leave it open.\n\n"
        "'Even the experts argue about this.' Leaving it unresolved is the design, not a gap."
    ))


def slide_question_card(prs):
    """Beat 7 — the unanswered question + two analogy keys."""
    s = add_blank(prs)
    add_title(s, "So… why bother learning anything?")

    # Hero question card — inverted (cream) so it pops on the dark slide
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              MARGIN, Inches(1.7), Inches(12.1), Inches(2.6))
    card.adjustments[0] = 0.04
    card.fill.solid(); card.fill.fore_color.rgb = HERO_FG
    card.line.fill.background()
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(40); tf.margin_right = Pt(40)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "If AI can do so much…\nwhy bother learning anything?"
    r.font.name = HEAD_FONT; r.font.size = Pt(44); r.font.color.rgb = BG
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(14)
    r2 = p2.add_run(); r2.text = "Don't answer. 13 weeks. On demo day — YOU tell US."
    r2.font.name = BODY_FONT; r2.font.size = Pt(14); r2.font.italic = True
    r2.font.color.rgb = RGBColor(0x55, 0x52, 0x4E)

    def key_card(x, w, eyebrow, title, body, color):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.6), w, Inches(2.1))
        c.adjustments[0] = 0.04
        c.fill.solid(); c.fill.fore_color.rgb = CARD_DARK
        c.line.color.rgb = color; c.line.width = Pt(1.5)
        c.shadow.inherit = False
        tf = c.text_frame
        tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(22); tf.margin_right = Pt(22)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = eyebrow.upper()
        r.font.name = BODY_FONT; r.font.size = Pt(10); r.font.bold = True
        r.font.color.rgb = color
        p2 = tf.add_paragraph(); p2.space_before = Pt(6)
        r2 = p2.add_run(); r2.text = title
        r2.font.name = HEAD_FONT; r2.font.size = Pt(19); r2.font.color.rgb = TEXT
        p3 = tf.add_paragraph(); p3.space_before = Pt(8)
        r3 = p3.add_run(); r3.text = body
        r3.font.name = BODY_FONT; r3.font.size = Pt(12); r3.font.color.rgb = MUTED

    half = (SLIDE_W - MARGIN * 3) / 2
    key_card(MARGIN, half,
             "Key #1 — calculator",
             "Why are you still learning times tables?",
             "Because the knowledge becomes a tool INSIDE YOUR OWN HEAD.",
             ACCENT)
    key_card(MARGIN * 2 + half, half,
             "Key #2 — car & exercise",
             "We've had cars 100 years. Did we stop running?",
             "No — moving your body is good for YOU. AI is the same for thinking.",
             P2)

    add_footer(s)
    add_notes(s, (
        "Plant the question UNANSWERED. She answers it on demo day. Hand her two keys to "
        "think with, as a MATCHED PAIR — let HER connect them to AI.\n\n"
        "KEY 1 — CALCULATOR: 'We've had calculators for 50 years — they're in every phone. "
        "So why are you still learning your times tables?' → because the knowledge becomes "
        "a tool INSIDE YOUR OWN HEAD.\n\n"
        "KEY 2 — CAR & EXERCISE: 'We've had cars for 100 years. Did we stop walking and "
        "running?' → no, because moving your body is good for YOU, however you get around. "
        "Honest wrinkle that makes it stronger: the car made it POSSIBLE to stop moving, so "
        "now we exercise on purpose. AI is the same for thinking — it makes it possible to "
        "stop, so learning becomes something we choose on purpose, because we know it's "
        "good for us.\n\n"
        "THE TURN: 'So — AI can write and draw and do math. What does that mean for "
        "learning to think?' Let her sit with it. Sign and date the question card with her; "
        "photograph it for page 1 of the exhibition site."
    ))


def slide_journal(prs):
    s = add_blank(prs)
    add_title(s, "Three I'd automate.  Three should stay human.")

    half = (SLIDE_W - MARGIN * 3) / 2

    def column(x, w, eyebrow, color):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, CONTENT_Y, w, CONTENT_H)
        c.adjustments[0] = 0.03
        c.fill.solid(); c.fill.fore_color.rgb = CARD_DARK
        c.line.color.rgb = color; c.line.width = Pt(1.5)
        c.shadow.inherit = False
        tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(28)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = eyebrow.upper()
        r.font.name = BODY_FONT; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = color
        for i in range(3):
            p = tf.add_paragraph(); p.space_before = Pt(16)
            r = p.add_run(); r.text = f"{i+1}.  ___________________________"
            r.font.name = BODY_FONT; r.font.size = Pt(17); r.font.color.rgb = TEXT

    column(MARGIN, half, "Automate", P2)
    column(MARGIN * 2 + half, half, "Keep human", P1)

    add_text(s, "Her words. The automate column will echo Beat 2 cards. Photograph → Page 1 of the exhibition site.",
             MARGIN, Inches(6.6), Inches(12.1), Inches(0.4),
             font=BODY_FONT, size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s)
    add_notes(s, (
        "Open the journal. Three in each column, in HER words. The 'automate' column will "
        "naturally echo the Beat 2 cards.\n\n"
        "Photograph the page → Page 1 of the live exhibition site."
    ))


# ---------------- Build ----------------

prs = new_pres()

# 1. Title
slide_title(prs)

# 2. Beat 1A — wow reel
video_slide(prs,
    "What can AI do right now?",
    "A self-driving Waymo, on a regular street",
    "Waymo driverless car ride POV 2025",
    "What you might cue up",
    [
        "Waymo driving itself.",
        "Humanoid robots dancing (Unitree, Spring Festival Gala).",
        "An AI-generated song — ideally live, from her own prompt.",
        "An AI-made image.",
    ],
    P1,
    ("Open by SHOWING, not defining. Roll a short WOW REEL of AI doing things humans do. "
     "Waymo is the primary clip; have 2–3 backups queued. Keep it short. Let the room "
     "react before you say anything analytical."),
)

# 3. Beat 1B — definition + milestones + provocation
video_slide(prs,
    "AI is mimicking us.",
    "AlphaGo's move 37 vs Lee Sedol",
    "AlphaGo move 37 Lee Sedol DeepMind documentary",
    "Where it has beaten the best of us",
    [
        "Chess — Deep Blue beat Kasparov, 1997.",
        "Go — AlphaGo beat Lee Sedol, 2016. Move 37: a move no human would play.",
        "Math — silver→gold *standard* at the international math olympiad (2024 → 2025).",
    ],
    P1,
    ("Land the definition: 'AI is when we build machines that mimic things humans normally "
     "do with their minds and senses — seeing, understanding language, deciding, creating.'\n\n"
     "For parents: deliberately broad and behavior-based. No single agreed technical "
     "definition — and that's part of the story.\n\n"
     "Quick milestone highlights (sidebar):\n"
     "  • CHESS: Deep Blue beat Kasparov, 1997.\n"
     "  • GO: AlphaGo beat Lee Sedol, 2016. Move 37.\n"
     "  • MATH: 2024 Google AI → silver-medal STANDARD at the IMO (28/42). 2025 Gemini "
     "Deep Think → gold-medal STANDARD (35/42), officially graded.\n\n"
     "HONESTY NOTE: say 'gold-medal STANDARD,' NEVER 'won gold' — the AI was not an "
     "official IMO contestant.\n\n"
     "Let ONE provocation hang: 'If AI can do ALL this… what are humans even for?'"),
    sidebar_footnote="'Gold-medal standard' — not an official contestant.",
)

# 4. Beat 2 — wishes (no video)
slide_beat2_wishes(prs)

# --- Beat 3 — backwards walk, each card reframed to 'what problem did this solve?' ---

video_slide(prs,
    "Video calling",
    "The Jetsons (1962) — videophone scene",
    "Jetsons 1962 videophone scene",
    "Problem it solved",
    [
        "Seeing the face of someone far away while you talk.",
        "1962: The Jetsons promised exactly this — wild future stuff.",
        "Now: in her pocket every week.",
    ],
    P3,
    ("Shreya video-calls family in Kolkata and Seattle. Ordinary. Sixty years ago — wild "
     "future. The story of every tool: yesterday's magic → today's normal."),
)

video_slide(prs,
    "The telephone",
    "Bell's 1876 telephone / early exchange archival",
    "Alexander Graham Bell telephone 1876 archival",
    "Problem it solved",
    [
        "Talking to someone far away — without writing it down.",
        "Bell, 1876.",
        "Before: letters that took weeks. Or shouting across a room.",
    ],
    P3,
    "Bell, 1876. First time a voice could leave a body and travel.",
)

video_slide(prs,
    "Recorded music",
    "Edison phonograph / gramophone playing",
    "Edison phonograph cylinder playing demonstration",
    "Problem it solved",
    [
        "Hearing a sound after it's already happened.",
        "Edison's phonograph, 1877.",
        "Before: somebody had to play it again. Every time.",
    ],
    P3,
    "A sound, frozen and replayable. The first 'pause / rewind' button in human history.",
)

video_slide(prs,
    "The camera",
    "First photograph + early box camera",
    "first photograph history Niepce Daguerre",
    "Problem it solved",
    [
        "Keeping what you saw — without painting it.",
        "Niépce, 1826. The first permanent photograph.",
        "Before: you hired an artist.",
    ],
    P3,
    "A moment, kept by anyone — not just the wealthy.",
)

video_slide(prs,
    "The printing press",
    "Working Gutenberg press demonstration",
    "Gutenberg printing press demonstration how it works",
    "Problem it solved",
    [
        "Copying a book — without a person writing it out by hand.",
        "Gutenberg, ~1440.",
        "Before: monks, months per copy.",
    ],
    P3,
    "Mass literacy. The world opens.",
)

video_slide(prs,
    "Writing",
    "Ancient cuneiform tablet / Egyptian hieroglyphs",
    "ancient cuneiform clay tablet writing history",
    "Problem it solved",
    [
        "Holding a thought outside your own head.",
        "~5000 years ago.",
        "Before: forget it, lose it forever.",
    ],
    P3,
    "The first time a thought could leave a brain and survive.",
)

video_slide(prs,
    "Fire & the first tools",
    "Cave painting + first fire / early humans",
    "early humans fire cave painting documentary",
    "Problem it solved",
    [
        "Making the world bend a little — to make life simpler.",
        "The first technology, ever.",
        "Cooked food → bigger brains. Warmth → winter. Light → night.",
    ],
    P3,
    "Every story starts here. The original tool.",
)

video_slide(prs,
    "The calculator",
    "Vintage mechanical + early electronic calculator demo",
    "vintage mechanical calculator and early electronic calculator demo",
    "Problem it solved  (we'll come back)",
    [
        "Math by hand: slow, error-prone.",
        "The most personal stop — Shreya is learning times tables right now.",
        "Plant the seed: we still learn them anyway. Why?",
    ],
    P3,
    ("The MOST PERSONAL stop. PLANT THE SEED — don't resolve it (Beat 7 is the payoff): "
     "times tables and the number line aren't only facts; they're THINKING-TOOLS. When you "
     "learn them, you install a tool INSIDE your own head that lets you see patterns and "
     "catch wrong answers.\n\n"
     "Background (presenter): Judy Fan's 'cognitive tools' idea. Shreya already uses one "
     "every week — written music notation captures sound that vanishes the instant it's "
     "made, and holds it still so she can think about it, share it, improve it."),
)

# --- Break (verbal) ---

# --- Beat 4 — what problems is AI solving NOW ---

video_slide(prs,
    "🐋  Talking to whales",
    "Project CETI — sperm whale communication",
    "Project CETI sperm whale communication AI",
    "What it does — and doesn't",
    [
        "Decoding sperm whale clicks using AI.",
        "Researchers BELIEVE they've found a 'phonetic alphabet.'",
        "We can't talk to whales yet — we're BEGINNING to understand them.",
        "Callback to Beat 2: if she wished she could 'talk to animals'…",
    ],
    P2,
    ("Project CETI uses AI to decode sperm whale clicks. Researchers have found what they "
     "BELIEVE is a phonetic alphabet, and even vowel-like sounds similar to human speech.\n\n"
     "HONESTY HEDGE: we CANNOT talk to whales yet — we're just BEGINNING to understand "
     "them. That gap is what makes it thrilling.\n\n"
     "CALLBACK to Beat 2's 'what do you wish you could do?' — if she said 'talk to "
     "animals,' this is the payoff.\n\n"
     "Sources: Project CETI; MIT CSAIL / Nature Communications (2024); Harvard SEAS (2025)."),
    sidebar_footnote="Sources: Project CETI; MIT CSAIL / Nature Comm (2024); Harvard SEAS (2025).",
)

video_slide(prs,
    "👓  Helping a blind person see",
    "Be My Eyes / Meta smart-glasses demo",
    "Be My Eyes Meta smart glasses blind assistance demo",
    "What it does",
    [
        "AI glasses describe what's in front of you — out loud.",
        "Reads menus, identifies objects and colors.",
        "Or connects to a human helper who sees through the camera.",
        "Still gets things wrong: noodles → 'lasagna'.",
    ],
    P2,
    ("AI glasses let a blind person ask 'describe what's in front of me' — the AI reads "
     "menus, identifies objects and colors, or connects them to a human volunteer (Be My "
     "Eyes) who sees through the camera and helps in real time.\n\n"
     "HONEST + FUNNY HEDGE: a blind user held up noodles, asked the glasses what they "
     "were, and the AI confidently answered 'LASAGNA.' Solves a real problem AND still "
     "gets things wrong — perfect for a 10-year-old to see.\n\n"
     "Sources: Meta / Be My Eyes (2025); Consumer Reports; CBC (2025)."),
    sidebar_footnote="Sources: Meta / Be My Eyes (2025); Consumer Reports; CBC.",
)

video_slide(prs,
    "🤖  A robot helped do surgery",
    "Johns Hopkins SRT-H — autonomous surgery robot",
    "Johns Hopkins SRT-H autonomous surgery robot",
    "What it did — and didn't",
    [
        "A robot learned by WATCHING surgeons.",
        "Performed part of a gallbladder operation.",
        "On a lifelike model — NOT a living person.",
        "Frame it as helping doctors, not replacing them.",
    ],
    P2,
    ("Johns Hopkins SRT-H learned to perform part of a gallbladder operation by WATCHING "
     "videos of surgeons, then carried it out and adapted in real time.\n\n"
     "CRUCIAL HEDGE: it was done on a LIFELIKE MODEL, not a living person. Frame as "
     "'helping doctors,' NOT replacing them.\n\n"
     "Source: Johns Hopkins / Science Robotics, July 2025."),
    sidebar_footnote="Source: Johns Hopkins / Science Robotics, July 2025.",
)

# --- Beat 5 — why copy us ---

video_slide(prs,
    "Experiences can't be outsourced",
    "Warm montage — dance, piano, swim, family meal",
    "Bharatanatyam performance children short",
    "AI can describe — never have",
    [
        "Goosebumps from a piece of music.",
        "The first taste of something delicious.",
        "Landing a hard Bharatanatyam step.",
        "Winning a swim race after weeks of trying.",
    ],
    P1,
    ("First answer to the provocation: EXPERIENCES.\n\n"
     "AI can DESCRIBE every one of these. It can never HAVE them.\n\n"
     "Ask Shreya: name one feeling AI could never feel."),
)

video_slide(prs,
    "Children are the great explorers",
    "Alison Gopnik — children, learning, exploration",
    "Alison Gopnik children learning exploration TED",
    "Why copy us?",
    [
        "Children are 'the best learning machines in the universe.' — Alison Gopnik",
        "Kids EXPLORE (wild ideas, why? why? why?). Adults EXPLOIT (narrow focus).",
        "Curiosity isn't a phase to grow out of — it's the superpower.",
    ],
    P1,
    ("Why mimic HUMANS, and not something else? Because the best learners and "
     "problem-solvers we know are human — and ESPECIALLY CHILDREN.\n\n"
     "Alison Gopnik calls children 'the best learning machines in the universe.' Kids are "
     "EXPLORERS — they experiment, ask why, try wild ideas ('explore' mode). Adults narrow "
     "into 'exploit' mode. She describes a child's attention as a broad LANTERN that "
     "lights up everything; an adult's is a narrow SPOTLIGHT.\n\n"
     "Make the implicit point land: curiosity isn't a phase to grow out of — it's "
     "Shreya's SUPERPOWER, and it's what AI is trying to copy TOWARD, not away from.\n\n"
     "Bonus (Gopnik 2025): today's AI models are 'cultural technologies' like the printing "
     "press or libraries — powerful ways to pass on what humans already know — NOT curious "
     "minds like a child.\n\n"
     "Sources: Gopnik, 'The Philosophical Baby'; 'Large AI models are cultural and social "
     "technologies,' Science (2025); Berkeley talks (2025–26)."),
    sidebar_footnote="Sources: Gopnik, 'The Philosophical Baby'; Science (2025).",
)

# Beat 5c — open question
slide_intelligence_open(prs)

# --- Beat 6 — what AI still can't do (with Beat 2 callback) ---

video_slide(prs,
    "The problems AI still CAN'T solve",
    "Home robot fumbling laundry — fail reel",
    "Weave Isaac laundry robot folding 2026",
    "What AI still can't do",
    [
        "Fold laundry reliably (Weave Isaac 0, ~$8K — still needs human help).",
        "Common sense about the physical world.",
        "Know what anything FEELS like.",
        "Remember your wishes from Beat 2? Those problems are YOURS.",
    ],
    P2,
    ("Ground the hope in real 2026 limits:\n"
     "  • Weave 'Isaac 0' laundry robot — ~$8K, STILL needs a human to remote-fix tricky "
     "folds.\n"
     "  • 1X 'NEO' home robot — leans on remote helpers for tasks it doesn't know.\n\n"
     "From the scientists:\n"
     "  • AI doesn't keep learning after it's built the way a kid does.\n"
     "  • It doesn't know what anything FEELS like.\n"
     "  • It lacks real common sense about the physical world.\n\n"
     "(Sources: LeCun–Malik–Dupoux 2026; Surya Ganguli, Daedalus 2026.)\n\n"
     "★ CALLBACK to Beat 2 — pick up her cards: 'Remember what you wished you didn't have "
     "to do — clean your room, fold the laundry? The most advanced robot on Earth STILL "
     "can't reliably do it. THAT PROBLEM IS YOURS.'\n\n"
     "Frame as INVITATION: so much is wide open for HER generation."),
    sidebar_footnote="LeCun–Malik–Dupoux (2026); Ganguli, Daedalus (2026).",
)

# Beat 7 — the question + analogy keys
slide_question_card(prs)

# Wrap
slide_journal(prs)


out = Path(__file__).parent / "deck.pptx"
prs.save(str(out))
print(f"Wrote {out}  ({len(prs.slides)} slides)")
